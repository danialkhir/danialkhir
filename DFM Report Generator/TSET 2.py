import sys
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QTabWidget, QLabel, QComboBox, QFileDialog, QScrollArea, QMessageBox
from PyQt5.QtGui import QPixmap, QColor
from pptx import Presentation
from pptx.util import Inches

class PowerPointGenerator(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.slides = []
        self.slide_count = 0
        self.slides_widgets = {} 
        self.finding_proposals = {
    'No IPC Class Level Stated': ['Proposal: Customer to provide IPC Class Level of Product'],
    'PCBA is a product component': ['Customer to provide AVL', 'Internal:Require EPA Assembly'],
    'Cable is a product component':['Customer to provide AVl','Intenal:IQC Check to Cater to IPC LEvel'],
    'Product require functional test':['Customer to Provide Test Fixture Setup','UWC to Produce Test Fixture for Customer'],
    'Product require Hi-Pot test':['Customer to define Hi-Pot Test Specification and Setup'],
    'Product require Continuity test':['Customer to provide complete schematic and continuity test point',"UWC to decide on the Continuity Test Point and to be verified by Customer Side"],
    'Product involve Spftware/Firmware/Hardware for Testing/Shipping':['Customer to provide required Software/Firmware/Hardware as listed in BOM'],
    'Product uses Custom part from supplier':['Customer to provide AVL','Internal:TO compare technical data sheet provided by AVL with incumbent'],
    'Product involves Electrical part to be attached on Mechanical Part':['Internal:NPI to be more Aware of the tolerance/dimension at Mechanical Attach Point of Electrical Component'],
    'Problem with Voltage Distribution Analysis':['Customer to Advise the Unmatched Voltage Level and/or Type'],
    'Problem with Power Load Analysis':['Customer to Advise Power Supply and Load Non-Balance Issue'],
    'Connector/Crimp used for Cable does not match mating part':['Customer to advice on the Actual Connector/Crimp required','UWC To Propose Suitable Crimp/Connector based on the Connection Point'],
    # proposal Content
    }

    def init_ui(self):
        self.setWindowTitle('DFM Report Generator')

        layout = QVBoxLayout()
        self.setStyleSheet("background-color: cyan; color: black;")  # background and warna

        tab_widget = QTabWidget()

        slide_tab = QWidget()
        self.slide_scroll = QScrollArea()
        self.slide_scroll.setWidgetResizable(True)
        self.slide_scroll_widget = QWidget()
        self.slide_layout = QVBoxLayout(self.slide_scroll_widget)
        self.slide_scroll.setWidget(self.slide_scroll_widget)

        self.add_slide_button = QPushButton('Add Part')
        self.add_slide_button.clicked.connect(self.add_slide)
        self.add_slide_button.setStyleSheet("background-color: yellow; color: black;")  # Setting button color
        self.slide_layout.addWidget(self.add_slide_button)

        self.preview_button = QPushButton('Preview Proposal')
        self.preview_button.clicked.connect(self.preview_slide)
        self.preview_button.setStyleSheet("background-color: yellow; color: black;")  # Setting button color
        self.slide_layout.addWidget(self.preview_button)

        self.publish_button = QPushButton('Publish Slide')
        self.publish_button.clicked.connect(self.publish_slide)
        self.publish_button.setStyleSheet("background-color: yellow; color: black;")  # Setting button color
        self.slide_layout.addWidget(self.publish_button)
        tab_widget.addTab(self.slide_scroll, 'Add Part')
        layout.addWidget(tab_widget)
        self.setLayout(layout)
        self.show()

    def add_slide(self):
        self.slide_count += 1

        slide_label = QLabel(f"Part {self.slide_count}")
        self.slide_layout.addWidget(slide_label)

        finding_combo = QComboBox()
        finding_combo.addItem('Select a finding')
        finding_combo.addItem('No IPC Class Level Stated')
        finding_combo.addItem('PCBA is a product component')
        finding_combo.addItem('Cable is a product component')
        finding_combo.addItem('Product require functional test')
        finding_combo.addItem('Product require Hi-Pot test')
        finding_combo.addItem('Product require Continuity test')
        finding_combo.addItem('Product involve Spftware/Firmware/Hardware for Testing/Shipping')
        finding_combo.addItem('Product uses Custom part from supplier')
        finding_combo.addItem('Product involves Electrical part to be attached on Mechanical Part')
        finding_combo.addItem('Problem with Voltage Distribution Analysis')
        finding_combo.addItem('Problem with Power Load Analysis')
        finding_combo.addItem('Connector/Crimp used for Cable does not match mating part')
        finding_combo.currentIndexChanged.connect(lambda: self.show_proposals(finding_combo.currentText()))
        self.slide_layout.addWidget(finding_combo)

        proposal_combo = QComboBox()
        proposal_combo.addItem('Select a proposal')
        proposal_combo.setObjectName(f"proposal_combo_{self.slide_count}")
        self.slide_layout.addWidget(proposal_combo)

        attach_picture_button = QPushButton(f'Attach Picture to Slide {self.slide_count}')
        attach_picture_button.clicked.connect(lambda: self.attach_picture(self.slide_count))  # Connect button to attach_picture method
        attach_picture_button.setStyleSheet("background-color: yellow; color: black;")  # Setting button color
        self.slide_layout.addWidget(attach_picture_button)

        picture_preview = QLabel()  # preview gambar
        picture_preview.setFixedSize(200, 200)
        self.slide_layout.addWidget(picture_preview)
        self.slides_widgets[self.slide_count] = {
        'slide_label': slide_label,
        'finding_combo': finding_combo,
        'proposal_combo': proposal_combo,
        'picture_preview': picture_preview
        }

    def show_proposals(self, finding):
        proposal_combo = self.sender().parent().findChild(QComboBox, f"proposal_combo_{self.slide_count}")
        proposal_combo.clear()
        proposal_combo.addItem('Select a proposal')
        proposals = self.finding_proposals.get(finding, [])
        proposal_combo.addItems(proposals)
    
    def preview_slide(self):
        current_slide = self.slides_widgets[self.slide_count]
        finding = current_slide['finding_combo'].currentText()
        proposal = current_slide['proposal_combo'].currentText()

        if finding != 'Select a finding' and proposal != 'Select a proposal':
            QMessageBox.information(self, "Slide Preview", f"Finding: {finding}\nProposal: {proposal}")
        else:
            QMessageBox.warning(self, "Incomplete Slide", "Please select a finding and a proposal.")

    def publish_slide(self):
        current_slide = self.slides_widgets[self.slide_count]
        finding = current_slide['finding_combo'].currentText()
        proposal = current_slide['proposal_combo'].currentText()

        if finding != 'Select a finding' and proposal != 'Select a proposal':
            self.generate_pptx()  # Pause here
            QMessageBox.information(self, "Slide Published", f"Slide with Finding: {finding}\nProposal: {proposal} has been published.")
        else:
            QMessageBox.warning(self, "Incomplete Slide", "Please select a finding and a proposal.")
    def generate_pptx(self):
       options = QFileDialog.Options()
       options |= QFileDialog.DontUseNativeDialog
       file_dialog = QFileDialog()
       file_dialog.setOptions(options)
       file_dialog.setFileMode(QFileDialog.AnyFile)
       file_dialog.setNameFilter("PowerPoint Files (*.pptx)")
       file_dialog.setDefaultSuffix("pptx")
       file_dialog.setAcceptMode(QFileDialog.AcceptSave)
       if file_dialog.exec_():
            file_paths = file_dialog.selectedFiles()
            if file_paths:
                file_path = file_paths[0]
                prs = Presentation()

                # Iterate through the slide_layout and gather information
                i = 0
                while i < self.slide_layout.count():
                    widget = self.slide_layout.itemAt(i).widget()

                    # Check if the widget is QLabel (for picture)
                    if isinstance(widget, QLabel):
                        # Logic to retrieve image path from QLabel
                        image_path = self.retrieve_image_path(widget)  # Implement this function

                        # Create a slide in the presentation and add the image
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        left = Inches(1)
                        top = Inches(1)
                        slide.shapes.add_picture(image_path, left, top, width=Inches(2), height=Inches(2))
                        i += 2  # Increment by 2 to skip the QComboBox

                    elif isinstance(widget, QComboBox) and widget.objectName().startswith("proposal_combo_"):
                        finding_combo = self.slide_layout.itemAt(i).widget()
                        finding = finding_combo.currentText()
                        proposal = widget.currentText()

                        # Use the 'finding' and 'proposal' variables to add content to the slide
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        title = slide.shapes.title
                        title.text = f"Slide {i // 3 + 1}: {finding}"
                        content = slide.placeholders[1]  # Placeholder index may vary
                        content.text = f"Proposal: {proposal}"
                        i += 1  # Move to the next QLabel widget

                    else:
                        i += 1  # Move to the next widget

                # Save the presentation
                prs.save(file_path)
                QMessageBox.information(self, "Presentation Created", f"PowerPoint presentation created successfully and saved at: {file_path}")

    def attach_picture(self, slide_number):
        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        file_name, _ = QFileDialog.getOpenFileName(self, "Select Image", "", "Image Files (*.png *.jpg *.jpeg *.gif)", options=options)

        if file_name:
            current_slide = self.slide_layout.itemAt(self.slide_layout.count() - 5).widget()

            picture_preview = QLabel()  # Create a QLabel for the picture preview
            picture_preview.setFixedSize(200, 200)
            pixmap = QPixmap(file_name)
            picture_preview.setPixmap(pixmap.scaled(200, 200, aspectRatioMode=True))

        # Set the image path as a tooltip for the QLabel
            picture_preview.setToolTip(file_name)

            self.slide_layout.insertWidget(self.slide_layout.count() - 1, picture_preview)  # Insert picture preview before the attachment button
        else:
            QMessageBox.warning(self, "Error", "No image selected.")

    def retrieve_image_path(self, label_widget):
    # Retrieve the image path from the QLabel's tooltip
        image_path = label_widget.toolTip()
        return image_path
def run_app():
    app = QApplication(sys.argv)
    ex = PowerPointGenerator()
    sys.exit(app.exec_())

if __name__ == '__main__':
    run_app()
