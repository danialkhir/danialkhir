import sys
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QLineEdit, QFileDialog, QMessageBox
from pptx import Presentation

class PowerPointGenerator(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()

    def init_ui(self):
        self.setWindowTitle('Text to PowerPoint Generator')

        layout = QVBoxLayout()

        self.text_input = QLineEdit()
        layout.addWidget(self.text_input)

        self.text_input2 = QLineEdit()
        layout.addWidget(self.text_input2)

        self.generate_button = QPushButton('Generate PowerPoint')
        self.generate_button.clicked.connect(self.generate_pptx)
        layout.addWidget(self.generate_button)

        self.setLayout(layout)
        self.show()

    def generate_pptx(self):
        text1 = self.text_input.text()
        text2 = self.text_input2.text()

        if text1.strip() == '' and text2.strip() == '':
            QMessageBox.warning(self, "Empty Input", "Please input some text.")
            return

        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        file_dialog = QFileDialog()
        file_dialog.setOptions(options)
        file_dialog.setFileMode(QFileDialog.AnyFile)
        file_dialog.setNameFilter("PowerPoint Files (*.pptx)")
        file_dialog.setDefaultSuffix("pptx")
        file_dialog.setAcceptMode(QFileDialog.AcceptSave)

        if file_dialog.exec_():
            file_paths = file_dialog.selectedFiles()
            if file_paths:
                file_path = file_paths[0]
                prs = Presentation()

                # Slide 1 with Title and Content (Index 0)
                if text1.strip() != '':
                    slide_layout = prs.slide_layouts[0]
                    slide1 = prs.slides.add_slide(slide_layout)
                    title1 = slide1.shapes.title
                    title1.text = "Title and Content Slide"
                    content1 = slide1.placeholders[1]
                    content1.text = text1

                # Slide 2 with Two Content Items (Index 3: Two Content)
                if text2.strip() != '':
                    slide_layout = prs.slide_layouts[3]
                    slide2 = prs.slides.add_slide(slide_layout)
                    title2 = slide2.shapes.title
                    title2.text = "Two Content Slide"
                    content2 = slide2.placeholders[0]  # Change the placeholder index based on the layout
                    content2.text = "Content 1"
                    content3 = slide2.placeholders[1]  # Change the placeholder index based on the layout
                    content3.text = text2

                prs.save(file_path)
                QMessageBox.information(self, "Presentation Created", f"PowerPoint presentation created successfully and saved at: {file_path}")

def run_app():
    app = QApplication(sys.argv)
    ex = PowerPointGenerator()
    sys.exit(app.exec_())

if __name__ == '__main__':
    run_app()
